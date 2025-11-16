"""
Script to generate a professional PowerPoint presentation for Road Safety Estimator
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime

def create_presentation():
    # Create a presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    PRIMARY_COLOR = RGBColor(30, 64, 175)  # Blue
    ACCENT_COLOR = RGBColor(102, 126, 234)  # Light Blue
    TEXT_COLOR = RGBColor(51, 51, 51)  # Dark Gray
    WHITE = RGBColor(255, 255, 255)
    
    # ===== SLIDE 1: WELCOME =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background gradient effect (using shape)
    background = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.gradient()
    background.fill.gradient_angle = 45
    background.fill.gradient_stops[0].color.rgb = RGBColor(102, 126, 234)
    background.fill.gradient_stops[1].color.rgb = RGBColor(118, 75, 162)
    background.line.fill.background()
    
    # Title
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🚦 ROAD SAFETY ESTIMATOR"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Cost Estimation for Road Safety Projects"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide1.shapes.add_textbox(
        Inches(1), Inches(5.5),
        Inches(8), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_frame.text = datetime.datetime.now().strftime("%B %Y")
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = WHITE
    date_para.alignment = PP_ALIGN.CENTER
    
    # ===== SLIDE 2: PROBLEM STATEMENT =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "The Problem"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add underline shape
    line = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    # Content
    content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    problems = [
        ("Time-Consuming Process", "Manual analysis takes hours to days per project"),
        ("Human Error", "Calculations prone to mistakes and inconsistencies"),
        ("Expert Knowledge Required", "Deep understanding of IRC standards needed"),
        ("Inconsistent Pricing", "Different estimators arrive at different costs"),
        ("Limited Scalability", "Cannot handle multiple projects efficiently")
    ]
    
    for i, (title, desc) in enumerate(problems):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = f"❌ {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.space_after = Pt(6)
        
        p2 = content_frame.add_paragraph()
        p2.text = f"   {desc}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_COLOR
        p2.space_after = Pt(18)
    
    # ===== SLIDE 3: OUR SOLUTION =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Our Solution"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Add underline
    line = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(2.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    # Subtitle
    subtitle_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Automated AI-Powered Cost Estimation System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.italic = True
    subtitle_para.font.color.rgb = ACCENT_COLOR
    
    # Key highlights in boxes
    highlights = [
        ("⚡ 95% Faster", "Minutes instead of hours"),
        ("🎯 99% Accurate", "Fuzzy matching with IRC standards"),
        ("📍 36 Locations", "All Indian states/UTs supported"),
        ("💰 Smart Pricing", "Location + inflation adjustments"),
        ("📊 Professional Reports", "PDF with IRC citations"),
        ("✉️ Email Integration", "Instant report sharing")
    ]
    
    y_pos = 2.3
    for i in range(0, len(highlights), 2):
        for j in range(2):
            if i + j < len(highlights):
                emoji, desc = highlights[i + j]
                x_pos = 0.8 + (j * 4.5)
                
                # Box background
                box = slide3.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x_pos), Inches(y_pos),
                    Inches(4), Inches(1)
                )
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(240, 245, 255)
                box.line.color.rgb = ACCENT_COLOR
                box.line.width = Pt(2)
                
                # Text
                text_box = slide3.shapes.add_textbox(
                    Inches(x_pos + 0.2), Inches(y_pos + 0.15),
                    Inches(3.6), Inches(0.7)
                )
                text_frame = text_box.text_frame
                text_frame.text = emoji
                p1 = text_frame.paragraphs[0]
                p1.font.size = Pt(18)
                p1.font.bold = True
                p1.font.color.rgb = PRIMARY_COLOR
                
                p2 = text_frame.add_paragraph()
                p2.text = desc
                p2.font.size = Pt(14)
                p2.font.color.rgb = TEXT_COLOR
        
        y_pos += 1.3
    
    # ===== SLIDE 4: HOW IT WORKS =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "How It Works"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Add underline
    line = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(2.8), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    # Workflow steps
    steps = [
        ("1. Upload Document", "PDF, DOCX, or TXT road safety audit report"),
        ("2. Intelligent Extraction", "AI identifies interventions, locations, quantities"),
        ("3. IRC Matching", "Fuzzy algorithm matches with IRC standards"),
        ("4. Cost Calculation", "Location + inflation + GST adjustments"),
        ("5. Generate Report", "Professional PDF with detailed breakdown"),
        ("6. Share via Email", "One-click delivery to stakeholders")
    ]
    
    y_start = 1.8
    for i, (step, desc) in enumerate(steps):
        # Step circle
        circle = slide4.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.7), Inches(y_start + (i * 0.85)),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY_COLOR
        circle.line.fill.background()
        
        # Step number
        num_box = slide4.shapes.add_textbox(
            Inches(0.7), Inches(y_start + (i * 0.85)),
            Inches(0.6), Inches(0.6)
        )
        num_frame = num_box.text_frame
        num_frame.text = str(i + 1)
        num_para = num_frame.paragraphs[0]
        num_para.font.size = Pt(24)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE
        num_para.alignment = PP_ALIGN.CENTER
        num_frame.vertical_anchor = 1  # Middle
        
        # Step text
        text_box = slide4.shapes.add_textbox(
            Inches(1.6), Inches(y_start + (i * 0.85)),
            Inches(7.5), Inches(0.7)
        )
        text_frame = text_box.text_frame
        text_frame.text = step
        p1 = text_frame.paragraphs[0]
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_COLOR
        
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_COLOR
        
        # Connector arrow
        if i < len(steps) - 1:
            arrow = slide4.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(0.85), Inches(y_start + (i * 0.85) + 0.65),
                Inches(0.3), Inches(0.15)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_COLOR
            arrow.line.fill.background()
    
    # ===== SLIDE 5: TECHNICAL ARCHITECTURE =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Technical Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Add underline
    line = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(4.2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    # Left column - Components
    comp_title = slide5.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.5), Inches(0.5))
    comp_frame = comp_title.text_frame
    comp_frame.text = "Core Components"
    comp_para = comp_frame.paragraphs[0]
    comp_para.font.size = Pt(24)
    comp_para.font.bold = True
    comp_para.font.color.rgb = PRIMARY_COLOR
    
    components = [
        ("📄 Document Parser", "PyPDF2, pdfplumber, python-docx"),
        ("🔍 Matching Engine", "Fuzzy matching, Levenshtein distance"),
        ("💰 Price Fetcher", "Location factors, inflation, GST"),
        ("📊 Report Generator", "FPDF, professional layouts"),
        ("✉️ Email Service", "Gmail SMTP, HTML emails")
    ]
    
    y_pos = 2.3
    for title, tech in components:
        text_box = slide5.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(4), Inches(0.7))
        text_frame = text_box.text_frame
        text_frame.text = title
        p1 = text_frame.paragraphs[0]
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_COLOR
        
        p2 = text_frame.add_paragraph()
        p2.text = tech
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_COLOR
        
        y_pos += 0.85
    
    # Right column - Tech Stack
    tech_title = slide5.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.5), Inches(0.5))
    tech_frame = tech_title.text_frame
    tech_frame.text = "Technology Stack"
    tech_para = tech_frame.paragraphs[0]
    tech_para.font.size = Pt(24)
    tech_para.font.bold = True
    tech_para.font.color.rgb = PRIMARY_COLOR
    
    tech_stack = [
        "🐍 Python 3.13.2",
        "🌐 Streamlit 1.51.0",
        "📊 Pandas, NumPy",
        "📄 PyPDF2, pdfplumber",
        "🔍 FuzzyWuzzy",
        "📈 FPDF, Plotly",
        "💾 Excel Database",
        "📧 SMTP Integration"
    ]
    
    y_pos = 2.3
    for item in tech_stack:
        text_box = slide5.shapes.add_textbox(Inches(5.4), Inches(y_pos), Inches(4), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = item
        p = text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        
        y_pos += 0.5
    
    # ===== SLIDE 6: KEY ACHIEVEMENTS & BENEFITS =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Key Achievements"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Add underline
    line = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(3.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    # Big stats
    stats = [
        ("95%", "Time Reduction", "Hours → Minutes"),
        ("99%", "Accuracy Rate", "IRC Matching"),
        ("36", "Locations", "States/UTs Supported")
    ]
    
    x_start = 0.8
    for stat, title, desc in stats:
        # Stat box
        box = slide6.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_start), Inches(1.8),
            Inches(2.8), Inches(1.5)
        )
        box.fill.gradient()
        box.fill.gradient_angle = 45
        box.fill.gradient_stops[0].color.rgb = PRIMARY_COLOR
        box.fill.gradient_stops[1].color.rgb = ACCENT_COLOR
        box.line.fill.background()
        
        # Stat number
        num_box = slide6.shapes.add_textbox(
            Inches(x_start), Inches(2),
            Inches(2.8), Inches(0.6)
        )
        num_frame = num_box.text_frame
        num_frame.text = stat
        num_para = num_frame.paragraphs[0]
        num_para.font.size = Pt(48)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE
        num_para.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide6.shapes.add_textbox(
            Inches(x_start), Inches(2.6),
            Inches(2.8), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide6.shapes.add_textbox(
            Inches(x_start), Inches(2.95),
            Inches(2.8), Inches(0.3)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = WHITE
        desc_para.alignment = PP_ALIGN.CENTER
        
        x_start += 3.1
    
    # Benefits list
    benefits_title = slide6.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.5))
    benefits_frame = benefits_title.text_frame
    benefits_frame.text = "Business Benefits"
    benefits_para = benefits_frame.paragraphs[0]
    benefits_para.font.size = Pt(24)
    benefits_para.font.bold = True
    benefits_para.font.color.rgb = PRIMARY_COLOR
    
    benefits = [
        "✅ Reduced operational costs and time",
        "✅ Consistent and transparent pricing",
        "✅ Professional documentation ready for stakeholders",
        "✅ Scalable for multiple concurrent projects",
        "✅ Easy to use - no IRC expertise required"
    ]
    
    y_pos = 4.3
    for benefit in benefits:
        text_box = slide6.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.5), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = benefit
        p = text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        
        y_pos += 0.5
    
    # ===== SLIDE 7: THANK YOU =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background gradient
    background = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.gradient()
    background.fill.gradient_angle = 45
    background.fill.gradient_stops[0].color.rgb = RGBColor(102, 126, 234)
    background.fill.gradient_stops[1].color.rgb = RGBColor(118, 75, 162)
    background.line.fill.background()
    
    # Thank you text
    thank_you_box = slide7.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.5)
    )
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(60)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = WHITE
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide7.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions & Discussion"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide7.shapes.add_textbox(
        Inches(2), Inches(5.5),
        Inches(6), Inches(1)
    )
    contact_frame = contact_box.text_frame
    contact_frame.text = "🚦 Road Safety Estimator\n📧 Streamlit Web Application\n🌐 Python-Based Solution"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        para.space_after = Pt(8)
    
    # Save presentation
    output_path = 'Road_Safety_Estimator_Presentation.pptx'
    prs.save(output_path)
    print(f"\nPresentation created successfully: {output_path}")
    print(f"Total slides: 7")
    print(f"  1. Welcome")
    print(f"  2. Problem Statement")
    print(f"  3. Our Solution")
    print(f"  4. How It Works")
    print(f"  5. Technical Architecture")
    print(f"  6. Key Achievements")
    print(f"  7. Thank You")
    
    return output_path

if __name__ == '__main__':
    create_presentation()
